from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os

# Folder with your plots
img_dir = "/work/gcelotto/btv_mini_rerun/CMSSW_10_6_32_patch1/src/plotsPerEvent"   # folder with your images
images = [f for f in os.listdir(img_dir) if f.lower().endswith((".png", ".jpg", ".jpeg"))]
images.sort()

prs = Presentation()
#blank_layout = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]
for idx,img_name in enumerate(images):
    #slide = prs.slides.add_slide(blank_layout)
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = f"Ev: {idx}"
    
    # Open image to get original size
    img_path = os.path.join(img_dir, img_name)
    with Image.open(img_path) as im:
        width_px, height_px = im.size

    # Slide dimensions in EMUs
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Compute scaling factor to fit slide while keeping aspect ratio
    scale_w = slide_w / width_px
    scale_h = slide_h / height_px
    scale = min(scale_w, scale_h)

    # Compute final width/height in EMUs
    final_w = int(width_px * scale)
    final_h = int(height_px * scale)

    # Center the image on slide
    left = int((slide_w - final_w) / 2)
    top = int((slide_h - final_h) / 2)

    slide.shapes.add_picture(img_path, left, top, width=final_w, height=final_h)

prs.save("/work/gcelotto/btv_mini_rerun/CMSSW_10_6_32_patch1/src/plotsPerEvent/plots_slides.pptx")

